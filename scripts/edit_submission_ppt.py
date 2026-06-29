"""Fill ISRO BAH 2026 Idea Submission Template — Bindas Code (formatted v2)."""
from __future__ import annotations

from pathlib import Path

import numpy as np
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch, FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ── paths ──────────────────────────────────────────────────────────────────
TEMPLATE = Path(
    r"c:\Users\Brindha\Downloads\ISRO Hackathon Clouds\[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx"
)
FALLBACK = Path(
    r"c:\Users\Brindha\Downloads\ISRO Hackathon Clouds\Bindas_Code_ISRO_BAH_2026_Submission_v5.pptx"
)
OUTPUT = Path(
    r"c:\Users\Brindha\Downloads\ISRO Hackathon Clouds\Bindas_Code_ISRO_BAH_2026_Submission_v6.pptx"
)
ASSETS = Path(r"c:\Users\Brindha\Downloads\ISRO Hackathon Clouds\outputs\ppt_assets")


def _open_presentation() -> Presentation:
    for candidate in (FALLBACK, TEMPLATE):
        if candidate.exists():
            prs = Presentation(str(candidate))
            if len(prs.slides) >= 10:
                return prs
    raise FileNotFoundError("No valid 10-slide template found.")

# ── team info ──────────────────────────────────────────────────────────────
TEAM_NAME = "Bindas Code"
TEAM_LEADER = "Brindha Manickavasakan"
COLLEGE = "Sri Ramakrishna Engineering College"
PROBLEM = "PS-2: Generative AI-Based Cloud Removal and Reconstruction for LISS-IV Satellite Imagery"
SOLUTION = "DCMF-UNet — Dual-Encoder SAR-Optical Cross-Modal Fusion GAN"
SUBTITLE = "SAR-guided cloud-free LISS-IV reconstruction for North Eastern Region (NER) India"

# ── reference palette (matches architect-playbook style) ───────────────────
C_BG = "#F5F4F0"
C_SAGE = "#6B8F71"
C_SAGE_DK = "#3D5A45"
C_TERRA = "#C4785A"
C_TERRA_DK = "#A8614A"
C_GRAY = "#8A8A8A"
C_DARK = "#2D3436"
C_WHITE = "#FFFFFF"
C_LIGHT = "#FAFAF8"


# ═══════════════════════════════════════════════════════════════════════════
#  Diagram helpers
# ═══════════════════════════════════════════════════════════════════════════


def _dot_grid(ax, xlim, ylim, spacing=0.25, color="#D8D8D4", alpha=0.6):
    for x in np.arange(xlim[0], xlim[1], spacing):
        for y in np.arange(ylim[0], ylim[1], spacing):
            ax.plot(x, y, ".", color=color, markersize=1.2, alpha=alpha)


def _rbox(ax, x, y, w, h, label, fc=C_WHITE, ec=C_SAGE_DK, fs=7.5, bold=False, lw=1.4):
    p = FancyBboxPatch(
        (x, y), w, h, boxstyle="round,pad=0.05,rounding_size=0.08",
        facecolor=fc, edgecolor=ec, linewidth=lw,
    )
    ax.add_patch(p)
    ax.text(x + w / 2, y + h / 2, label, ha="center", va="center",
            fontsize=fs, color=C_DARK, fontweight="bold" if bold else "normal",
            wrap=True, linespacing=1.3)
    return p


def _arrow(ax, x1, y1, x2, y2, color=C_TERRA, lw=2.0):
    ax.add_patch(FancyArrowPatch(
        (x1, y1), (x2, y2), arrowstyle="-|>", mutation_scale=14,
        color=color, linewidth=lw, shrinkA=2, shrinkB=2,
    ))


def _callout(ax, tx, ty, bx, by, text, side="left"):
    ax.annotate(
        text, xy=(bx, by), xytext=(tx, ty),
        fontsize=7, color=C_GRAY, fontstyle="italic",
        arrowprops=dict(arrowstyle="-|>", color=C_GRAY, lw=1.0,
                        connectionstyle="arc3,rad=0.25"),
        ha="center" if side == "center" else ("right" if side == "left" else "left"),
    )


def _step_card(
    ax, x, y, w, h, title, details,
    header_h=0.72, ec=C_SAGE_DK, header_fc=C_TERRA, header_ec=C_TERRA_DK,
    title_fs=8.5, detail_fs=6.5,
):
    """Single process-step card: one outer border, coloured header band, white body."""
    lw = 2.0
    r = 0.14

    # Outer card — single clean border
    outer = FancyBboxPatch(
        (x, y), w, h,
        boxstyle=f"round,pad=0,rounding_size={r}",
        facecolor=C_WHITE, edgecolor=ec, linewidth=lw, zorder=2,
    )
    ax.add_patch(outer)

    # Header band (inset so it stays inside the outer border)
    pad = lw * 0.004 + 0.04
    hdr = FancyBboxPatch(
        (x + pad, y + h - header_h - pad), w - 2 * pad, header_h,
        boxstyle=f"round,pad=0,rounding_size={r * 0.85}",
        facecolor=header_fc, edgecolor=header_ec, linewidth=1.2, zorder=3,
    )
    ax.add_patch(hdr)

    # Divider between header and body
    div_y = y + h - header_h - pad
    ax.plot(
        [x + pad + 0.06, x + w - pad - 0.06], [div_y, div_y],
        color=header_ec, linewidth=1.0, zorder=4, solid_capstyle="round",
    )

    # Text
    ax.text(
        x + w / 2, y + h - header_h / 2 - pad / 2, title,
        ha="center", va="center", fontsize=title_fs, fontweight="bold",
        color=C_WHITE, zorder=5, linespacing=1.2,
    )
    body_cy = y + (h - header_h - pad) / 2
    ax.text(
        x + w / 2, body_cy, details,
        ha="center", va="center", fontsize=detail_fs,
        color=C_DARK, zorder=5, linespacing=1.35,
    )


def _layer_box(ax, x, y, w, h, title, inner_pad=0.28):
    """Container with a clean double-line border and title bar."""
    # Outer border
    outer = FancyBboxPatch(
        (x, y), w, h,
        boxstyle="round,pad=0,rounding_size=0.12",
        facecolor=C_WHITE, edgecolor=C_SAGE_DK, linewidth=2.2, zorder=1,
    )
    ax.add_patch(outer)
    # Inner accent line (inset)
    inset = 0.06
    inner = FancyBboxPatch(
        (x + inset, y + inset), w - 2 * inset, h - 2 * inset,
        boxstyle="round,pad=0,rounding_size=0.10",
        facecolor="none", edgecolor=C_SAGE, linewidth=0.8, linestyle="--", zorder=1,
    )
    ax.add_patch(inner)
    ax.text(
        x + w / 2, y + h - inner_pad, title,
        ha="center", va="center", fontsize=9, fontweight="bold", color=C_SAGE_DK, zorder=2,
    )
    return x + inset + 0.12, y + inset + 0.12, w - 2 * inset - 0.24, h - inner_pad - inset - 0.18


def _inner_tab(ax, x, y, w, h, label, active=False, fs=6.5):
    """Tab box that sits cleanly inside a layer container."""
    fc = C_SAGE if active else C_WHITE
    ec = C_SAGE_DK
    text_color = C_WHITE if active else C_DARK
    tab = FancyBboxPatch(
        (x, y), w, h,
        boxstyle="round,pad=0,rounding_size=0.08",
        facecolor=fc, edgecolor=ec, linewidth=1.6, zorder=3,
    )
    ax.add_patch(tab)
    ax.text(
        x + w / 2, y + h / 2, label,
        ha="center", va="center", fontsize=fs, color=text_color,
        fontweight="bold" if active else "normal", zorder=4, linespacing=1.25,
    )


def draw_system_architecture(path: Path) -> None:
    """Full system architecture — spaced layout, no overlapping text."""
    fig, ax = plt.subplots(figsize=(16, 9))
    ax.set_xlim(0, 16)
    ax.set_ylim(0, 9)
    ax.axis("off")
    ax.set_facecolor(C_BG)
    _dot_grid(ax, (0, 16), (0, 9))

    ax.text(
        8, 8.55, "System Architecture — LISS-IV Cloud Removal Platform",
        ha="center", fontsize=14, fontweight="bold", color=C_DARK,
    )

    # ── Left: user entry flow ─────────────────────────────────────────────
    _rbox(ax, 0.35, 6.85, 1.35, 0.72, "Analyst /\nISRO User", fc=C_LIGHT, ec=C_TERRA, bold=True, fs=8)
    _arrow(ax, 1.72, 7.21, 2.05, 7.21)
    _rbox(ax, 2.05, 6.62, 1.75, 1.05, "Streamlit\nDashboard\n(Pattern Router)",
          fc=C_SAGE, ec=C_SAGE_DK, fs=7.5, bold=True)

    _arrow(ax, 3.82, 7.45, 4.35, 7.75)
    _arrow(ax, 3.82, 6.95, 4.35, 6.55)
    _rbox(ax, 4.35, 7.55, 1.45, 0.62, "Real-time\nUpload & Infer", fc=C_LIGHT, ec=C_TERRA, fs=7)
    _rbox(ax, 4.35, 6.45, 1.45, 0.62, "Batch\nDownload & Train", fc=C_LIGHT, ec=C_TERRA, fs=7)

    # ── Execution Layer (wide container) ──────────────────────────────────
    ex, ey, ew, eh = 5.85, 3.65, 7.55, 3.05
    _layer_box(ax, ex, ey, ew, eh, "Execution Layer")

    # Sub-label — inside top-left, below title bar
    ax.text(ex + 0.35, ey + eh - 0.62, "Data Pipeline Tools",
            fontsize=7.5, fontweight="bold", color=C_TERRA_DK, va="top")

    # Four pipeline tools — fit left zone only, stop before guardrails column
    tool_labels = ["Bhoonidhi\nDownload", "Band Stack\nG/R/NIR", "Cloud Mask\nGen.", "SAR\nCo-register"]
    tg, ty, th = 0.16, ey + 1.05, 1.05
    tx0 = ex + 0.35
    gx = ex + ew - 2.05
    tool_zone_w = gx - tx0 - 0.18
    tw = (tool_zone_w - tg * (len(tool_labels) - 1)) / len(tool_labels)
    for i, lbl in enumerate(tool_labels):
        _rbox(ax, tx0 + i * (tw + tg), ty, tw, th, lbl, fc=C_LIGHT, ec=C_TERRA, fs=6.5)

    # Guardrails panel — right column inside execution layer
    gy, gw, gh = ty, 1.75, th
    lock = FancyBboxPatch(
        (gx, gy), 0.38, gh, boxstyle="round,pad=0.03",
        facecolor=C_SAGE_DK, edgecolor=C_SAGE_DK, linewidth=1.2, zorder=3,
    )
    ax.add_patch(lock)
    ax.text(gx + 0.19, gy + gh / 2, "G", ha="center", va="center",
            fontsize=9, color="white", fontweight="bold", zorder=4)
    _rbox(
        ax, gx + 0.45, gy, gw - 0.45, gh,
        "Quality Guardrails\nCRS / GeoTIFF check\nClear-pixel preserve\nMask / spectral checks",
        fc=C_SAGE, ec=C_SAGE_DK, fs=5.8,
    )
    _callout(ax, gx - 0.1, ey + eh + 0.35, gx + gw / 2, ey + eh + 0.05,
             "Application intercepts\nguarding the core", "center")

    # Model bar — bottom strip inside execution layer, below tools
    model_y = ey + 0.28
    _rbox(
        ax, ex + 0.35, model_y, ew - 0.7, 0.62,
        "DCMF-UNet GAN  |  Optical + SAR Encoders + Fusion Gates + Decoder\n"
        "(CrossModalFusionGate opens toward SAR inside clouds)",
        fc=C_TERRA, ec=C_TERRA_DK, fs=5.8, bold=True,
    )

    # Arrows into execution layer
    _arrow(ax, 5.82, 7.15, 5.85, 5.55, color=C_SAGE)
    _arrow(ax, 5.82, 6.75, 5.85, 4.85, color=C_SAGE)

    # ── Synthesis output — fully outside execution layer ──────────────────
    sx, sy, sw, sh = 14.05, 4.55, 1.75, 1.85
    _arrow(ax, ex + ew + 0.02, 5.55, sx - 0.05, 5.55, color=C_TERRA, lw=2.5)
    _rbox(
        ax, sx, sy, sw, sh,
        "Synthesis\n\nCloud-free GeoTIFF\nSide-by-side PNG\nPSNR / SSIM / SAM\nNDVI consistency",
        fc=C_LIGHT, ec=C_TERRA_DK, fs=6.5, bold=True,
    )

    # ── State Management — bottom, centred, below execution ───────────────
    _arrow(ax, ex + ew / 2, ey - 0.02, ex + ew / 2, 3.38, color=C_TERRA)
    sx0, sy0, sw0, sh0 = 3.15, 0.55, 10.2, 2.65
    _layer_box(ax, sx0, sy0, sw0, sh0,
               "State Management  (Checkpoints + Patch Store + Geo Context)")

    bw, bg = 2.95, 0.22
    by = sy0 + 0.45
    bx0 = sx0 + 0.35
    _rbox(ax, bx0, by, bw, 1.65,
          "Pruning Logic\n\nPatch overlap merge\nFeather blend at edges\nMask dilation control",
          fc=C_TERRA, ec=C_TERRA_DK, fs=6.2)
    _rbox(ax, bx0 + bw + bg, by, bw, 1.65,
          "Shared Patch Store\n\ndata/fusion/ pairs\nlissiv_cloudy / clear\nsentinel1 / masks",
          fc=C_SAGE, ec=C_SAGE_DK, fs=6.2)
    _rbox(ax, bx0 + 2 * (bw + bg), by, 2.65, 1.65,
          "Context Window\n\n256 x 256 patches\n50% overlap\nGeoTIFF CRS kept",
          fc=C_LIGHT, ec=C_SAGE_DK, fs=6.2)
    _callout(ax, sx0 + sw0 / 2, sy0 - 0.38, sx0 + sw0 / 2, sy0 + 0.02,
             "Shared memory sustaining the lifecycle", "center")

    # ── HALO companion — bottom-left, clear of state box ──────────────────
    _rbox(
        ax, 0.25, 0.65, 2.65, 2.35,
        "HALO Companion\n(Hybrid Atmospheric-Latent\nOptical Reconstruction)\n\n"
        "Physics inversion for thin cloud\nGenerative fill for opaque holes\nSelf-supervised via CloudForge",
        fc=C_LIGHT, ec=C_GRAY, fs=6.2,
    )

    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor=C_BG, pad_inches=0.12)
    plt.close(fig)


def draw_model_architecture(path: Path) -> None:
    """DCMF-UNet internal architecture — reference style."""
    fig, ax = plt.subplots(figsize=(14, 8))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 8)
    ax.axis("off")
    ax.set_facecolor(C_BG)
    _dot_grid(ax, (0, 14), (0, 8))

    ax.text(7, 7.55, "Model Architecture — DCMF-UNet (Dual-Encoder Cross-Modal Fusion U-Net)",
            ha="center", fontsize=13, fontweight="bold", color=C_DARK)

    # Inputs
    _rbox(ax, 0.4, 5.5, 1.5, 0.75, "LISS-IV Cloudy\nG / R / NIR", fc=C_LIGHT, ec=C_TERRA, bold=True)
    _rbox(ax, 0.4, 4.4, 1.5, 0.75, "Cloud Mask\n(1 = reconstruct)", fc=C_LIGHT, ec=C_GRAY)
    _rbox(ax, 0.4, 3.0, 1.5, 0.75, "Sentinel-1 SAR\nVV + VH", fc=C_LIGHT, ec=C_SAGE, bold=True)
    _callout(ax, 1.15, 6.5, 1.15, 6.25, "Multi-modal\ninputs", "center")

    # Dual encoders layer
    _layer_box(ax, 2.3, 2.8, 5.5, 3.8, "Dual-Encoder Feature Extraction")
    _rbox(ax, 2.6, 5.0, 2.2, 1.3,
          "Optical Encoder\n(mask-aware gated convs)\nInput: optical + mask",
          fc=C_TERRA, ec=C_TERRA_DK, fs=7, bold=True)
    _rbox(ax, 2.6, 3.2, 2.2, 1.3,
          "SAR Encoder\n(conv blocks)\nSeparate radar statistics",
          fc=C_SAGE, ec=C_SAGE_DK, fs=7, bold=True)

    # Fusion gates at 3 scales
    ax.text(5.2, 6.85, "CrossModalFusionGate (×3 scales)", fontsize=7.5,
            fontweight="bold", color=C_SAGE_DK)
    for i, y in enumerate([5.5, 4.5, 3.5]):
        _rbox(ax, 5.0, y, 2.3, 0.75,
              f"Fusion Gate Scale {i+1}\ngate = f(optical, SAR, mask)\nfused = opt·(1-g) + merge·g",
              fc=C_LIGHT, ec=C_SAGE_DK, fs=6)
    _callout(ax, 6.15, 7.3, 6.15, 6.9, "Gate opens toward\nSAR inside clouds", "center")

    _arrow(ax, 1.9, 5.85, 2.6, 5.65)
    _arrow(ax, 1.9, 4.75, 2.6, 5.55)
    _arrow(ax, 1.9, 3.35, 2.6, 3.85)

    # Bottleneck + Decoder
    _layer_box(ax, 8.1, 3.5, 2.5, 3.1, "Reconstruction Head")
    _rbox(ax, 8.4, 5.5, 1.9, 0.8, "SE Bottleneck\n(global context)", fc=C_SAGE, ec=C_SAGE_DK, fs=7)
    _rbox(ax, 8.4, 4.2, 1.9, 1.0, "U-Net Decoder\n(gated conv + skip connections)", fc=C_TERRA, ec=C_TERRA_DK, fs=7)
    _arrow(ax, 7.3, 5.0, 8.4, 5.5)

    # Discriminator
    _rbox(ax, 8.4, 3.65, 1.9, 0.45, "PatchGAN 70×70 Discriminator", fc=C_LIGHT, ec=C_GRAY, fs=6.5)
    ax.text(9.35, 3.55, "Sees [optical | SAR] — enforces radar-consistent texture",
            ha="center", fontsize=6, color=C_GRAY, fontstyle="italic")

    # Output compositing
    _arrow(ax, 10.6, 4.8, 11.3, 4.8, color=C_TERRA, lw=2.5)
    _rbox(ax, 11.3, 4.0, 2.2, 1.6,
          "Cloud-Free Output\n\noutput = optical·(1−mask)\n       + generated·mask\n\nClear pixels preserved\nexactly — no hallucination",
          fc=C_LIGHT, ec=C_TERRA_DK, fs=7, bold=True)

    # Loss bar
    _rbox(ax, 2.3, 1.0, 11.0, 1.5,
          "Training Loss Function\n"
          "L = λ₁·Masked L1  +  λ₂·Perceptual (VGG16)  +  λ₃·LSGAN  +  λ₄·Spectral (SAM + band stats)  +  λ₅·(1 − SSIM)\n"
          "Generator ~8M params  |  Discriminator ~2.8M params  |  Patch size 256×256  |  Inference: tiled + feather blend",
          fc=C_WHITE, ec=C_SAGE_DK, fs=7)

    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor=C_BG)
    plt.close(fig)


def draw_process_flow(path: Path) -> None:
    """End-to-end process flow — unified step cards with clean borders."""
    fig, ax = plt.subplots(figsize=(14, 7.5))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 7.5)
    ax.axis("off")
    ax.set_facecolor(C_BG)
    _dot_grid(ax, (0, 14), (0, 7.5))

    ax.text(
        7, 7.05, "End-to-End Process Flow — LISS-IV Cloud Removal Pipeline",
        ha="center", fontsize=13, fontweight="bold", color=C_DARK,
    )

    steps = [
        ("1. Data\nAcquisition", "Bhoonidhi API\nLISS-IV MX70 L2\nNER 8 states"),
        ("2. Preprocess", "Stack BAND2/3/4\n→ G/R/NIR\nDownsample 2048 px"),
        ("3. Cloud Mask", "Heuristic / Fmask\nWhite = cloud pixel\nMask dilation"),
        ("4. SAR Fusion", "Sentinel-1 VV/VH\nCo-register to\n5.8 m grid"),
        ("5. Train GAN", "DCMF-UNet\n10 epochs\nColab GPU"),
        ("6. Deliver", "GeoTIFF export\nMetrics report\nPoster viz"),
    ]

    n = len(steps)
    gap = 0.32
    margin_l, margin_r = 0.45, 0.45
    usable = 14 - margin_l - margin_r
    card_w = (usable - gap * (n - 1)) / n
    card_h = 2.35
    card_y = 3.55

    for i, (title, detail) in enumerate(steps):
        cx = margin_l + i * (card_w + gap)
        _step_card(ax, cx, card_y, card_w, card_h, title, detail)
        if i < n - 1:
            ax_x1 = cx + card_w + 0.04
            ax_x2 = cx + card_w + gap - 0.04
            mid_y = card_y + card_h / 2
            _arrow(ax, ax_x1, mid_y, ax_x2, mid_y, color=C_TERRA, lw=2.2)

    # Streamlit dashboard layer — tabs inset inside container borders
    lx, ly, lw, lh = _layer_box(
        ax, 0.55, 0.55, 12.9, 2.15,
        "Streamlit Dashboard — User Interface Layer",
    )
    tabs = [
        ("Results\n(upload & infer)", True),
        ("Pipeline\n(overview)", False),
        ("Download\n(Bhoonidhi)", False),
        ("Preprocess\n(band stack)", False),
        ("Train\n(DCMF-UNet)", False),
    ]
    tab_gap = 0.18
    tab_w = (lw - tab_gap * (len(tabs) - 1)) / len(tabs)
    tab_h = lh - 0.05
    tab_y = ly
    for i, (label, active) in enumerate(tabs):
        _inner_tab(ax, lx + i * (tab_w + tab_gap), tab_y, tab_w, tab_h, label, active=active)

    _callout(ax, 7, 2.95, 7, 3.45, "Single dashboard\norchestrates full pipeline", "center")

    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor=C_BG, pad_inches=0.15)
    plt.close(fig)


def draw_ui_wireframe(path: Path) -> None:
    """Streamlit UI wireframe — reference style."""
    fig, ax = plt.subplots(figsize=(14, 7))
    ax.set_xlim(0, 14)
    ax.set_ylim(0, 7)
    ax.axis("off")
    ax.set_facecolor(C_BG)
    _dot_grid(ax, (0, 14), (0, 7))

    ax.text(7, 6.55, "Wireframe — Streamlit Cloud Removal Dashboard",
            ha="center", fontsize=13, fontweight="bold", color=C_DARK)

    # Sidebar
    sb = FancyBboxPatch((0.4, 0.5), 2.6, 5.8, boxstyle="round,pad=0.06",
                         facecolor=C_SAGE_DK, edgecolor=C_SAGE_DK, linewidth=2)
    ax.add_patch(sb)
    ax.text(1.7, 5.95, "Navigation", ha="center", color="white", fontsize=9, fontweight="bold")
    for i, item in enumerate(["Results", "Pipeline", "Download", "Preprocess", "Train"]):
        fc = C_TERRA if item == "Results" else C_SAGE
        _rbox(ax, 0.65, 4.9 - i * 0.85, 2.1, 0.65, item, fc=fc, ec=C_WHITE, fs=8, bold=(item == "Results"))

    # Main panel
    _layer_box(ax, 3.3, 0.5, 10.3, 5.8, "Results — Cloud Removal Output")
    _rbox(ax, 3.6, 5.3, 9.7, 0.55, "Upload cloudy LISS-IV image  (PNG / JPG / GeoTIFF)", fc=C_WHITE, ec=C_GRAY, fs=8)
    for x, lbl in [(3.8, "Cloud brightness"), (6.2, "Cloud whiteness"), (8.6, "Edge expand")]:
        ax.add_patch(mpatches.Rectangle((x, 4.55), 2.0, 0.12, facecolor="#E8E8E4", edgecolor=C_GRAY))
        ax.plot(x + 1.0, 4.61, "o", color=C_TERRA, markersize=8)
        ax.text(x + 1.0, 4.35, lbl, ha="center", fontsize=7, color=C_GRAY)

    _rbox(ax, 3.6, 1.0, 4.5, 3.1, "Cloudy LISS-IV\n(input)\n\n[Satellite scene\nwith cloud cover]", fc="#FFF5F0", ec=C_TERRA, fs=8)
    _rbox(ax, 8.5, 1.0, 4.5, 3.1, "Cloud-Free\nReconstruction\n\n[DCMF-UNet output\nG/R/NIR restored]", fc="#F0FFF4", ec=C_SAGE, fs=8)
    _arrow(ax, 8.1, 2.5, 8.5, 2.5, color=C_TERRA, lw=2.5)
    ax.text(8.3, 2.75, "Infer", ha="center", fontsize=7, color=C_TERRA, fontweight="bold")

    fig.savefig(path, dpi=200, bbox_inches="tight", facecolor=C_BG)
    plt.close(fig)


# ═══════════════════════════════════════════════════════════════════════════
#  PowerPoint helpers
# ═══════════════════════════════════════════════════════════════════════════

def _set_cell(cell, text: str, bold: bool = False, size: int = 11, align=PP_ALIGN.LEFT):
    cell.text = text
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for p in cell.text_frame.paragraphs:
        p.alignment = align
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.name = "Calibri"


def _style_table_header(table, row: int = 0, fill=RGBColor(0x3D, 0x5A, 0x45)):
    for ci in range(len(table.columns)):
        cell = table.cell(row, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.bold = True


def _fill_grouped_table(
    slide,
    groups: list[tuple[str, list[tuple[str, str]]]],
    headers: tuple[str, str, str],
    col_widths: tuple[float, float, float],
    *,
    top=Inches(0.95),
    height=Inches(4.4),
    group_fill=RGBColor(0xE8, 0xF0, 0xE8),
    stripe_fill=RGBColor(0xF5, 0xF4, 0xF0),
) -> None:
    """Build a 3-column table with merged group labels in column 0."""
    total_rows = 1 + sum(len(items) for _, items in groups)
    tbl = _add_table(slide, total_rows, 3, Inches(0.35), top, Inches(9.3), height)
    tbl.columns[0].width = Inches(col_widths[0])
    tbl.columns[1].width = Inches(col_widths[1])
    tbl.columns[2].width = Inches(col_widths[2])

    for ci, header in enumerate(headers):
        _set_cell(tbl.cell(0, ci), header, bold=True, size=10, align=PP_ALIGN.CENTER)
    _style_table_header(tbl)

    row = 1
    for gi, (group_name, items) in enumerate(groups):
        start_row = row
        for item_name, detail in items:
            _set_cell(tbl.cell(row, 1), item_name, bold=True, size=9)
            _set_cell(tbl.cell(row, 2), detail, size=8)
            if (row - start_row) % 2 == 1:
                for ci in (1, 2):
                    tbl.cell(row, ci).fill.solid()
                    tbl.cell(row, ci).fill.fore_color.rgb = stripe_fill
            row += 1

        group_cell = tbl.cell(start_row, 0)
        _set_cell(group_cell, group_name, bold=True, size=9, align=PP_ALIGN.CENTER)
        group_cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        if len(items) > 1:
            group_cell.merge(tbl.cell(row - 1, 0))
        group_cell.fill.solid()
        group_cell.fill.fore_color.rgb = group_fill
        for p in group_cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER


def _add_table(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, left, top, width, height).table


def _clear_text_shapes(slide, keep_bg=True):
    """Remove all shapes except background picture."""
    sp_tree = slide.shapes._spTree
    to_remove = []
    for shape in slide.shapes:
        if keep_bg and shape.shape_type == 13:  # PICTURE
            continue
        to_remove.append(shape)
    for shape in to_remove:
        sp_tree.remove(shape._element)


def _add_title(slide, text: str, top=Inches(0.35)):
    box = slide.shapes.add_textbox(Inches(0.4), top, Inches(9.2), Inches(0.55))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.name = "Calibri"
    return box


def _add_image(slide, path: Path, top, height=Inches(4.8)):
    slide.shapes.add_picture(str(path), Inches(0.35), top, width=Inches(9.3))


# ═══════════════════════════════════════════════════════════════════════════
#  Slide fillers
# ═══════════════════════════════════════════════════════════════════════════

def fill_slide1_cover(slide) -> None:
    """Structured cover slide with project title + info table."""
    _clear_text_shapes(slide)

    # Project title block
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.55), Inches(9.0), Inches(1.1))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Generative AI-Based Cloud Removal and Reconstruction"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "for LISS-IV Satellite Imagery"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = SUBTITLE
    p3.font.size = Pt(13)
    p3.font.name = "Calibri"
    p3.font.italic = True
    p3.alignment = PP_ALIGN.CENTER

    # Cover info table
    rows, cols = 8, 2
    tbl = _add_table(slide, rows, cols, Inches(0.7), Inches(2.0), Inches(8.6), Inches(2.6))
    tbl.columns[0].width = Inches(2.8)
    tbl.columns[1].width = Inches(5.8)

    cover_data = [
        ("Team Name", TEAM_NAME),
        ("Team Leader", TEAM_LEADER),
        ("College", COLLEGE),
        ("Problem Statement", PROBLEM),
        ("Proposed Solution", SOLUTION),
        ("Target Region", "North Eastern Region (NER) — Assam, Meghalaya, Nagaland, Manipur, Mizoram, Tripura, Arunachal Pradesh, Sikkim"),
        ("Hackathon", "ISRO Bharatiya Antariksh Hackathon (BAH) 2026"),
        ("Core Innovation", "SAR-optical fusion with mask-guided CrossModalFusionGate — preserves clear pixels, reconstructs only cloud regions"),
    ]
    for ri, (label, value) in enumerate(cover_data):
        _set_cell(tbl.cell(ri, 0), label, bold=True, size=11)
        _set_cell(tbl.cell(ri, 1), value, size=10)
        if ri % 2 == 0:
            for ci in range(2):
                tbl.cell(ri, ci).fill.solid()
                tbl.cell(ri, ci).fill.fore_color.rgb = RGBColor(0xF5, 0xF4, 0xF0)


def fill_slide2_team(slide) -> None:
    for shape in slide.shapes:
        if shape.has_table:
            t = shape.table
            t.cell(0, 0).text = f"Team Leader:\n\nName: {TEAM_LEADER}\nCollege: {COLLEGE}\nRole: ML / Full-stack Lead"


def fill_slide3_opportunity(slide) -> None:
    _clear_text_shapes(slide)
    _add_title(slide, "Opportunity — Problem, Differentiation & USP")

    tbl = _add_table(slide, 5, 2, Inches(0.4), Inches(1.05), Inches(9.2), Inches(4.3))
    tbl.columns[0].width = Inches(2.2)
    tbl.columns[1].width = Inches(7.0)
    _style_table_header(tbl)

    rows = [
        ("The Problem",
         "LISS-IV optical imagery over NER India is frequently unusable during monsoon (May–Sep) due to thick cloud cover. "
         "Clouds block Green/Red/NIR bands at ~5.8 m resolution, halting agriculture monitoring, disaster response, and land-use analysis. "
         "Optical-only inpainting fails where the surface is fully occluded."),
        ("How We Differ",
         "• Dual-modality fusion: LISS-IV optical + Sentinel-1 SAR (VV/VH) — SAR penetrates cloud and preserves terrain structure\n"
         "• DCMF-UNet with CrossModalFusionGate at every encoder scale — not a flat detect→fuse→reconstruct pipeline\n"
         "• Mask-aware gated convolutions + clear-pixel preservation: output = optical×(1−mask) + generated×mask\n"
         "• HALO companion: physics-based radiometric recovery (J = (I−A(1−τ))/τ) for thin cloud; generative fill only for opaque holes\n"
         "• Self-supervised CloudForge training when paired cloudy/clear data is scarce"),
        ("How It Solves",
         "End-to-end platform: Bhoonidhi download → band stack (BAND2/3/4) → cloud mask → SAR co-registration → "
         "GAN training on 256×256 patches → tiled GeoTIFF inference with CRS preserved → PSNR/SSIM/SAM evaluation. "
         "Streamlit dashboard enables non-expert analysts to run the full pipeline. Google Colab notebook for GPU training."),
        ("USP",
         "• SAR-guided reconstruction anchored to real geometry — not hallucinated texture\n"
         "• Zero hallucination on clear land — only cloud pixels are synthesized\n"
         "• GeoTIFF in/out with full geospatial metadata (CRS, transform)\n"
         "• Metrics aligned to ISRO evaluation: PSNR > 32 dB, SSIM > 0.90, SAM < 5°, NDVI RMSE < 0.05\n"
         "• Works with as few as 50–100 paired patches; synthetic bootstrap available"),
    ]
    headers = ["Aspect", "Details"]
    _set_cell(tbl.cell(0, 0), headers[0], bold=True, size=11)
    _set_cell(tbl.cell(0, 1), headers[1], bold=True, size=11)
    for ri, (aspect, detail) in enumerate(rows, start=1):
        _set_cell(tbl.cell(ri, 0), aspect, bold=True, size=10)
        _set_cell(tbl.cell(ri, 1), detail, size=9)


def fill_slide4_features(slide) -> None:
    _clear_text_shapes(slide)
    _add_title(slide, "Features Offered by the Solution")

    groups = [
        (
            "Data Acquisition\n& Preprocessing",
            [
                ("Bhoonidhi API Integration",
                 "Search & batch-download LISS-IV MX70 L2 for all 8 NER states; STAC catalog; IP-whitelisted access"),
                ("Automated Preprocessing",
                 "Stack BAND2/3/4 → G/R/NIR; downsample to 2048 px (~5 MB/scene); full-res option (~3.5 GB/scene)"),
                ("Cloud Mask Generation",
                 "Heuristic RGB brightness/whiteness; Fmask / s2cloudless support; adjustable dilation"),
                ("SAR Co-registration",
                 "Sentinel-1 VV/VH aligned to LISS-IV 5.8 m grid via fusion/preprocess/coregister.py"),
            ],
        ),
        (
            "AI / ML Core",
            [
                ("DCMF-UNet Generator",
                 "Dual-encoder cross-modal fusion U-Net (~8M params); 3-scale CrossModalFusionGate"),
                ("PatchGAN Discriminator",
                 "Conditional 70×70 PatchGAN (~2.8M params); sees [optical|SAR] for radar-consistent texture"),
                ("Multi-Loss Training",
                 "Masked L1 + VGG perceptual + LSGAN + spectral (SAM) + SSIM; synthetic bootstrap mode"),
                ("HALO Framework",
                 "Self-supervised physics + generative pipeline; continuous opacity τ; uncertainty map"),
            ],
        ),
        (
            "Application\n& Deployment",
            [
                ("GeoTIFF Inference",
                 "Tiled processing, 50% overlap, feather blend; CRS and transform preserved"),
                ("Streamlit Dashboard",
                 "5-page app: Results, Pipeline, Download, Preprocess, Train — live demo ready"),
                ("Google Colab Training",
                 "Self-contained colab_training/ folder with GPU notebook"),
            ],
        ),
        (
            "Quality\nAssurance",
            [
                ("Evaluation Suite",
                 "PSNR, SSIM, SAM on full image and cloud-mask regions; NDVI consistency check"),
            ],
        ),
    ]
    _fill_grouped_table(
        slide, groups,
        headers=("Feature Group", "Component", "Description"),
        col_widths=(1.55, 2.35, 5.4),
    )


def fill_slide8_technologies(slide) -> None:
    _clear_text_shapes(slide)
    _add_title(slide, "Technologies Used in the Solution")

    groups = [
        (
            "Core Stack",
            [
                ("Python 3.x", "Core language for entire pipeline"),
                ("PyTorch, TorchVision", "DCMF-UNet GAN training & inference"),
                ("NumPy, scikit-image", "Array operations; PSNR / SSIM / SAM metrics"),
            ],
        ),
        (
            "Geospatial\n& Satellite Data",
            [
                ("Rasterio, GDAL", "GeoTIFF read/write; CRS & transform preservation"),
                ("Bhoonidhi STAC API", "LISS-IV MX70 L2 download (Resourcesat-2/2A)"),
                ("LISS-IV MX70 L2", "Green, Red, NIR @ ~5.8 m — NER India study area"),
                ("Sentinel-1 (VV, VH)", "Cloud-penetrating SAR; ±6–12 day acquisition window"),
            ],
        ),
        (
            "Preprocessing\nPipelines",
            [
                ("Fmask, s2cloudless, heuristic", "Cloud mask generation — white = reconstruct"),
                ("coregister.py, patch_extract.py", "SAR alignment; 256×256 patch extraction"),
            ],
        ),
        (
            "Model\nArchitecture",
            [
                ("DCMF-UNet", "Dual-encoder, CrossModalFusionGate, SE bottleneck, U-Net decoder"),
                ("HALO + CloudForge", "Physics inversion; self-supervised cloud forging"),
            ],
        ),
        (
            "Web &\nDeployment",
            [
                ("Streamlit", "Interactive dashboard — download, train, visualize"),
                ("Google Colab", "GPU-accelerated training notebook"),
                ("Pillow (PIL)", "Upload preview, PNG export, poster visualizations"),
                ("Requests", "Bhoonidhi authentication & HTTP download"),
            ],
        ),
    ]
    _fill_grouped_table(
        slide, groups,
        headers=("Technology Group", "Tool / Library", "Purpose / Role"),
        col_widths=(1.55, 2.35, 5.4),
    )


def fill_slide9_cost(slide) -> None:
    _clear_text_shapes(slide)
    _add_title(slide, "Estimated Implementation Cost")

    cost_rows = [
        ("Bhoonidhi LISS-IV data", "₹0", "Free for registered ISRO Bhoonidhi users"),
        ("Sentinel-1 SAR (Copernicus)", "₹0", "Open access via Copernicus Hub / GEE"),
        ("Python + PyTorch + Rasterio", "₹0", "Open-source software stack"),
        ("Streamlit hosting", "₹0", "Free tier on Streamlit Community Cloud"),
        ("Development laptop", "₹0", "Already available (team hardware)"),
        ("Google Colab (GPU training)", "₹0 – ₹1,000 / month", "Free tier sufficient; Pro optional for faster training"),
        ("Cloud storage (~50 GB scenes)", "₹100 – ₹300 / month", "Optional — local disk used for prototype"),
        ("Domain / SSL certificate", "₹0 – ₹500 / year", "Optional for custom deployment URL"),
        ("Production GPU server", "₹5,000 – ₹12,000 / month", "Only if deploying beyond hackathon (AWS/GCP GPU instance)"),
        ("TOTAL — Prototype", "₹0 – ₹2,000", "Entire hackathon prototype on free/open-source tools"),
        ("TOTAL — Production", "₹5,000 – ₹15,000 / month", "GPU inference server + storage + hosting"),
    ]
    tbl = _add_table(slide, len(cost_rows) + 1, 3, Inches(0.35), Inches(0.95), Inches(9.3), Inches(4.4))
    tbl.columns[0].width = Inches(3.2)
    tbl.columns[1].width = Inches(2.0)
    tbl.columns[2].width = Inches(4.1)
    _style_table_header(tbl)
    _set_cell(tbl.cell(0, 0), "Component", bold=True, size=10)
    _set_cell(tbl.cell(0, 1), "Cost (INR)", bold=True, size=10)
    _set_cell(tbl.cell(0, 2), "Notes", bold=True, size=10)
    for ri, (comp, cost, notes) in enumerate(cost_rows, start=1):
        bold = comp.startswith("TOTAL")
        _set_cell(tbl.cell(ri, 0), comp, bold=bold, size=9 if not bold else 10)
        _set_cell(tbl.cell(ri, 1), cost, bold=bold, size=9 if not bold else 10, align=PP_ALIGN.CENTER)
        _set_cell(tbl.cell(ri, 2), notes, size=8)
        if bold:
            for ci in range(3):
                tbl.cell(ri, ci).fill.solid()
                tbl.cell(ri, ci).fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xE8)


def fill_slide10_thanks(slide) -> None:
    _clear_text_shapes(slide)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(4.0))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [
        ("Thank You", 28, True),
        ("", 8, False),
        (f"Team {TEAM_NAME}", 18, True),
        (COLLEGE, 14, False),
        ("", 8, False),
        ("Problem Statement: PS-2 — Cloud Removal from LISS-IV Satellite Imagery", 13, False),
        (SOLUTION, 13, True),
        ("", 8, False),
        (f"Team Leader: {TEAM_LEADER}", 14, False),
        ("", 8, False),
        ("Turning unusable cloudy LISS-IV scenes into analysis-ready, cloud-free imagery", 12, False),
        ("for agriculture, disaster response, and land monitoring across North East India.", 12, False),
        ("", 8, False),
        ("ISRO Bharatiya Antariksh Hackathon 2026", 14, True),
    ]
    for i, (text, size, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = "Calibri"


def fill_diagram_slides(prs: Presentation, assets: Path) -> None:
    assets.mkdir(parents=True, exist_ok=True)
    flow = assets / "process_flow_v2.png"
    system = assets / "system_architecture_v2.png"
    model = assets / "model_architecture_v2.png"
    ui = assets / "ui_wireframe_v2.png"

    draw_process_flow(flow)
    draw_system_architecture(system)
    draw_model_architecture(model)  # saved to assets for reference / appendix
    draw_ui_wireframe(ui)

    # Slide 5 — process flow
    s5 = prs.slides[4]
    _clear_text_shapes(s5)
    _add_title(s5, "Process Flow Diagram")
    _add_image(s5, flow, Inches(0.85))

    # Slide 6 — UI wireframe
    s6 = prs.slides[5]
    _clear_text_shapes(s6)
    _add_title(s6, "Wireframes / Mock UI — Streamlit Dashboard")
    _add_image(s6, ui, Inches(0.85))

    # Slide 7 — system architecture (reference style, full slide)
    s7 = prs.slides[6]
    _clear_text_shapes(s7)
    _add_title(s7, "Architecture Diagram — End-to-End Platform")
    _add_image(s7, system, Inches(0.85), height=Inches(4.6))


def main() -> None:
    prs = _open_presentation()
    fill_slide1_cover(prs.slides[0])
    fill_slide2_team(prs.slides[1])
    fill_slide3_opportunity(prs.slides[2])
    fill_slide4_features(prs.slides[3])
    fill_diagram_slides(prs, ASSETS)
    fill_slide8_technologies(prs.slides[7])
    fill_slide9_cost(prs.slides[8])
    fill_slide10_thanks(prs.slides[9])

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Saved: {OUTPUT}")
    try:
        prs.save(str(TEMPLATE))
        print(f"Updated template: {TEMPLATE}")
    except PermissionError:
        print("Template locked — close PowerPoint and re-run to update original.")


if __name__ == "__main__":
    main()
